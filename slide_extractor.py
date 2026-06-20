"""
Slide text extraction — supports PPTX and PDF.
Used by /api/upload-slides to extract per-slide content for relevance analysis.
"""
import os


def extract_slides(file_path: str) -> list[str]:
    """Return a list of text strings, one per slide/page (0-indexed)."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext in ('.pptx', '.ppt'):
        return _from_pptx(file_path)
    if ext == '.pdf':
        return _from_pdf(file_path)
    raise ValueError(f"Unsupported format '{ext}'. Use .pptx or .pdf")


def _from_pptx(path: str) -> list[str]:
    from pptx import Presentation
    prs = Presentation(path)
    result = []
    for slide in prs.slides:
        lines = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    lines.append(t)
        result.append('\n'.join(lines))
    return result


def _from_pdf(path: str) -> list[str]:
    import pypdf
    reader = pypdf.PdfReader(path)
    return [page.extract_text() or '' for page in reader.pages]
